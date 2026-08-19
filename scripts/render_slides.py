#!/usr/bin/env python3
"""Generate worship-slide PPTX files in image or editable render mode."""

import argparse
import json
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

W, H = 1920, 1080
LINE_GAP, PAD_H, PAD_V, BOX_RADIUS, SHADOW_OFFSET = 24, 80, 45, 18, 4
MAX_FONT_SIZE, MIN_FONT_SIZE, MAX_LINE_WIDTH = 90, 48, int(W * 0.78)
RES_MIN_W, RES_MIN_H, RES_WARN_W, RES_WARN_H, RES_DOWNSAMPLE_W = 1280, 720, 1920, 1080, 3840
DELIVERY_QUALITY = {"screen": 75, "standard": 92, "print": 98}
RENDER_MODES = ("image", "editable")
RENDER_MODE_KEYWORD_PREFIX = "church-worship-slides:render-mode="
# Editable-mode typography must survive font substitution on macOS/Keynote and Windows.
# "Arial" and "PingFang TC" ship with macOS; "Microsoft JhengHei" ships with Windows.
# Latin runs use Arial (metric-compatible with Liberation Sans, used for measurement).
EDITABLE_LATIN_FONT = "Arial"
EDITABLE_CJK_FONT = "PingFang TC"
# Font actually used to measure editable geometry, so the box matches the delivered font.
EDITABLE_METRIC_FONTS = [
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
]
# Substituted fonts can be wider/taller than the measured font; reserve headroom.
EDITABLE_WIDTH_SAFETY = 1.22
EDITABLE_LINE_HEIGHT_FACTOR = 1.32
EDITABLE_AUTOFIT_FONT_SCALE = 92500  # OOXML fontScale: allow ~7.5% shrink before overflow
EDITABLE_AUTOFIT_LINE_REDUCTION = 10000
FONT_CANDIDATES = [
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Black.ttc",
    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
]


def validate_background(bg_path):
    """Validate and preprocess the source background per the existing policy."""
    img = Image.open(bg_path).convert("RGBA")
    width, height = img.size
    if width < RES_MIN_W or height < RES_MIN_H:
        print(
            f"ERROR: Background image is too low-resolution ({width}x{height}). "
            f"Minimum required is {RES_MIN_W}x{RES_MIN_H}.",
            file=sys.stderr,
        )
        sys.exit(1)
    if width < RES_WARN_W or height < RES_WARN_H:
        print(f"WARNING: Background image ({width}x{height}) is below Full HD; it will be upscaled.")
    if width > RES_DOWNSAMPLE_W:
        new_height = int(height * RES_DOWNSAMPLE_W / width)
        img = img.resize((RES_DOWNSAMPLE_W, new_height), Image.Resampling.LANCZOS)
        print(f"INFO: Source image downsampled to {RES_DOWNSAMPLE_W}x{new_height} (exceeds 4K cap).")
    return img


def load_font(size):
    for path in FONT_CANDIDATES:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                pass
    return ImageFont.load_default()


def auto_font_size(lines, requested_size):
    """Use the established pixel-width guard to choose one size for the song."""
    size = requested_size if requested_size > 0 else MAX_FONT_SIZE
    draw = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    while size >= MIN_FONT_SIZE:
        font = load_font(size)
        widths = [draw.textbbox((0, 0), line, font=font)[2] for line in lines]
        if max(widths) <= MAX_LINE_WIDTH:
            return font, size
        size -= 4
    return load_font(MIN_FONT_SIZE), MIN_FONT_SIZE


def load_editable_metric_font(size):
    """Load the font whose metrics approximate what macOS/Windows will actually render."""
    for path in EDITABLE_METRIC_FONTS:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                pass
    return load_font(size)


def editable_auto_font_size(lines, requested_size):
    """Choose the largest size that fits the safe width using substitution-tolerant metrics."""
    if requested_size > 0:
        return requested_size
    draw = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    size = MAX_FONT_SIZE
    while size > MIN_FONT_SIZE:
        metric_font = load_editable_metric_font(size)
        widest = max(draw.textlength(line, font=metric_font) for line in lines)
        if widest * EDITABLE_WIDTH_SAFETY <= MAX_LINE_WIDTH:
            return size
        size -= 2
    return MIN_FONT_SIZE


def editable_layout_for_text(text, font_size_px):
    """Size the overlay box for PowerPoint/Keynote line boxes, not tight ink bounds.

    Uses a substitution-tolerant metric font, line-height based vertical sizing, and
    width/height safety headroom so a wider substituted font still fits inside the box.
    """
    lines = text.split("\n")
    metric_font = load_editable_metric_font(font_size_px)
    draw = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    widths = [draw.textlength(line, font=metric_font) for line in lines]
    text_width = max(widths) * EDITABLE_WIDTH_SAFETY
    line_height = font_size_px * EDITABLE_LINE_HEIGHT_FACTOR
    text_height = line_height * len(lines)

    box_w = int(min(W - 40, text_width + PAD_H * 2))
    box_h = int(text_height + PAD_V * 2)
    return {
        "lines": lines,
        "x": max(20, (W - box_w) // 2),
        "y": max(30, int(H / 4 - box_h / 2)),
        "w": box_w,
        "h": box_h,
    }


def line_metrics(lines, font):
    draw = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    boxes = [draw.textbbox((0, 0), line, font=font) for line in lines]
    metrics = [(x1 - x0, y1 - y0, y0) for x0, y0, x1, y1 in boxes]
    return max(item[0] for item in metrics), sum(item[1] for item in metrics) + LINE_GAP * (len(lines) - 1), metrics


def layout_for_text(text, font):
    """Return shared Full-HD geometry for review images and editable objects."""
    lines = text.split("\n")
    text_width, text_height, metrics = line_metrics(lines, font)
    box_w, box_h = text_width + PAD_H * 2, text_height + PAD_V * 2
    return {
        "lines": lines,
        "metrics": metrics,
        "x": max(20, (W - box_w) // 2),
        "y": max(30, int(H / 4 - box_h / 2)),
        "w": box_w,
        "h": box_h,
    }


def rounded_rectangle(draw, x0, y0, x1, y1, radius, fill):
    draw.rectangle([x0 + radius, y0, x1 - radius, y1], fill=fill)
    draw.rectangle([x0, y0 + radius, x1, y1 - radius], fill=fill)
    draw.ellipse([x0, y0, x0 + radius * 2, y0 + radius * 2], fill=fill)
    draw.ellipse([x1 - radius * 2, y0, x1, y0 + radius * 2], fill=fill)
    draw.ellipse([x0, y1 - radius * 2, x0 + radius * 2, y1], fill=fill)
    draw.ellipse([x1 - radius * 2, y1 - radius * 2, x1, y1], fill=fill)


def cover_background(bg_img):
    bg = bg_img.copy()
    if bg.width / bg.height > W / H:
        new_width = int(H * bg.width / bg.height)
        bg = bg.resize((new_width, H), Image.Resampling.LANCZOS)
        return bg.crop(((new_width - W) // 2, 0, (new_width - W) // 2 + W, H))
    new_height = int(W * bg.height / bg.width)
    bg = bg.resize((W, new_height), Image.Resampling.LANCZOS)
    return bg.crop((0, (new_height - H) // 2, W, (new_height - H) // 2 + H))


def render_slide(text, bg_img, font, overlay_rgb, overlay_alpha):
    """Create the review JPEG and the flattened-image slide canvas."""
    bg = Image.alpha_composite(cover_background(bg_img), Image.new("RGBA", (W, H), (0, 0, 0, 50)))
    layout = layout_for_text(text, font)
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    rounded_rectangle(
        ImageDraw.Draw(overlay), layout["x"], layout["y"], layout["x"] + layout["w"], layout["y"] + layout["h"],
        BOX_RADIUS, (*overlay_rgb, overlay_alpha),
    )
    bg = Image.alpha_composite(bg, overlay)
    text_layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(text_layer)
    current_y = layout["y"] + PAD_V
    for (ink_width, ink_height, y_offset), line in zip(layout["metrics"], layout["lines"]):
        x, y = (W - ink_width) // 2, current_y - y_offset
        draw.text((x + SHADOW_OFFSET, y + SHADOW_OFFSET), line, font=font, fill=(0, 0, 0, 200))
        draw.text((x, y), line, font=font, fill=(255, 255, 255, 255))
        current_y += ink_height + LINE_GAP
    return Image.alpha_composite(bg, text_layer).convert("RGB")


def px_to_emu(value, slide_extent, canvas_extent):
    return int(value / canvas_extent * slide_extent)


def set_fill(shape, rgb, opacity=255):
    """Apply a solid fill with an OOXML alpha value so it stays editable in PowerPoint."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    color = fill.fore_color._color._xClr
    alpha = color.find(qn("a:alpha"))
    if alpha is None:
        alpha = OxmlElement("a:alpha")
        color.append(alpha)
    alpha.set("val", str(round(max(0, min(255, opacity)) / 255 * 100000)))


def remove_outline(shape):
    shape.line.fill.background()


def set_portable_typefaces(run, latin_font, cjk_font):
    """Set Latin, East Asian, and complex-script typefaces to fonts present on macOS/Windows.

    Without an explicit East Asian typeface, Keynote/PowerPoint substitute an arbitrary
    font whose metrics differ from the measured layout, which caused v1.4.0 overflow.
    """
    rpr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", latin_font), ("a:ea", cjk_font), ("a:cs", cjk_font)):
        element = rpr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            rpr.append(element)
        element.set("typeface", typeface)


def enable_shrink_on_overflow(text_frame, font_scale, line_reduction):
    """Emit OOXML normAutofit so renderers shrink text instead of spilling out of the box."""
    body_pr = text_frame._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        existing = body_pr.find(qn(tag))
        if existing is not None:
            body_pr.remove(existing)
    autofit = OxmlElement("a:normAutofit")
    autofit.set("fontScale", str(font_scale))
    autofit.set("lnSpcReduction", str(line_reduction))
    body_pr.append(autofit)


def add_cover_background(slide, bg_path, slide_width, slide_height):
    """Place the original image as a cover-cropped native PowerPoint picture."""
    with Image.open(bg_path) as image:
        image_ratio = image.width / image.height
    slide_ratio = slide_width / slide_height
    picture = slide.shapes.add_picture(bg_path, 0, 0, width=slide_width, height=slide_height)
    if image_ratio > slide_ratio:
        crop = (1 - slide_ratio / image_ratio) / 2
        picture.crop_left = picture.crop_right = crop
    elif image_ratio < slide_ratio:
        crop = (1 - image_ratio / slide_ratio) / 2
        picture.crop_top = picture.crop_bottom = crop


def add_editable_slide(prs, text, bg_path, font, font_size_px, overlay_rgb, overlay_alpha):
    """Add an editable background picture, dimmer, rounded overlay, and lyric textbox."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_cover_background(slide, bg_path, prs.slide_width, prs.slide_height)

    dimmer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(dimmer, (0, 0, 0), 50)
    remove_outline(dimmer)

    layout = editable_layout_for_text(text, font_size_px)
    x = px_to_emu(layout["x"], prs.slide_width, W)
    y = px_to_emu(layout["y"], prs.slide_height, H)
    width = px_to_emu(layout["w"], prs.slide_width, W)
    height = px_to_emu(layout["h"], prs.slide_height, H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    set_fill(overlay, overlay_rgb, overlay_alpha)
    remove_outline(overlay)

    # The text box spans the overlay exactly; padding lives in the box geometry, so the
    # text frame keeps only a small inset and relies on autofit for the final guard.
    textbox = slide.shapes.add_textbox(x, y, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    inset_h = px_to_emu(PAD_H // 3, prs.slide_width, W)
    inset_v = px_to_emu(PAD_V // 3, prs.slide_height, H)
    frame.margin_left = frame.margin_right = inset_h
    frame.margin_top = frame.margin_bottom = inset_v
    enable_shrink_on_overflow(frame, EDITABLE_AUTOFIT_FONT_SCALE, EDITABLE_AUTOFIT_LINE_REDUCTION)

    font_size_pt = Pt(font_size_px * 72 / 96)
    for index, line in enumerate(layout["lines"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = EDITABLE_LINE_HEIGHT_FACTOR
        for run in paragraph.runs:
            run.font.name = EDITABLE_LATIN_FONT
            run.font.size = font_size_pt
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            set_portable_typefaces(run, EDITABLE_LATIN_FONT, EDITABLE_CJK_FONT)


def set_presentation_render_mode(prs, render_mode):
    props = prs.core_properties
    previous = [item.strip() for item in (props.keywords or "").split(";") if item.strip() and not item.strip().startswith(RENDER_MODE_KEYWORD_PREFIX)]
    props.keywords = "; ".join(previous + [f"{RENDER_MODE_KEYWORD_PREFIX}{render_mode}"])
    props.subject = "Church worship lyric slides"


def get_presentation_render_mode(prs):
    """Return the recorded mode, or None for a pre-v1.4.0 deck without a marker."""
    for item in (prs.core_properties.keywords or "").split(";"):
        item = item.strip()
        if item.startswith(RENDER_MODE_KEYWORD_PREFIX):
            mode = item.removeprefix(RENDER_MODE_KEYWORD_PREFIX)
            return mode if mode in RENDER_MODES else None
    return None


def create_pptx(lyrics_data, bg_path, output_path, overlay_rgb, overlay_alpha, font_size, jpeg_quality, render_mode):
    bg_img = validate_background(bg_path)
    all_lines = [line for slide in lyrics_data for line in slide["text"].split("\n")]
    font, used_size = auto_font_size(all_lines, font_size)
    # Editable decks are measured with a substitution-tolerant font, so they get their own size.
    editable_size = editable_auto_font_size(all_lines, font_size) if render_mode == "editable" else used_size
    reported_size = editable_size if render_mode == "editable" else used_size
    print(f"Render: {W}x{H} | mode: {render_mode} | font: {reported_size}px | overlay: rgb{overlay_rgb} alpha={overlay_alpha}")
    preview_dir = os.path.join(os.path.dirname(os.path.abspath(output_path)), "_slide_imgs")
    os.makedirs(preview_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    set_presentation_render_mode(prs, render_mode)
    for slide_data in lyrics_data:
        preview_path = os.path.join(preview_dir, f"slide_{slide_data['index']:02d}.jpg")
        render_slide(slide_data["text"], bg_img, font, overlay_rgb, overlay_alpha).save(preview_path, "JPEG", quality=jpeg_quality, optimize=True)
        if render_mode == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(preview_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            add_editable_slide(prs, slide_data["text"], bg_path, font, editable_size, overlay_rgb, overlay_alpha)
        print(f"  Slide {slide_data['index']:02d}: preview {os.path.getsize(preview_path) // 1024} KB")
    prs.save(output_path)
    print(f"Saved: {output_path} ({os.path.getsize(output_path) // 1024} KB)")


def parse_color(value):
    try:
        color = tuple(int(item.strip()) for item in value.split(","))
    except ValueError as exc:
        raise argparse.ArgumentTypeError("Colour must be three comma-separated integers.") from exc
    if len(color) != 3 or any(item < 0 or item > 255 for item in color):
        raise argparse.ArgumentTypeError("Colour values must be integers from 0 to 255.")
    return color


def main():
    parser = argparse.ArgumentParser(description="Generate worship slides PPTX")
    parser.add_argument("--lyrics", required=True)
    parser.add_argument("--bg", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--overlay-color", default="0,0,0")
    parser.add_argument("--overlay-alpha", type=int, default=185)
    parser.add_argument("--font-size", type=int, default=0)
    parser.add_argument("--delivery-mode", default="standard", choices=list(DELIVERY_QUALITY))
    parser.add_argument("--render-mode", default="image", choices=RENDER_MODES, help="image = flattened; editable = native PowerPoint shapes")
    args = parser.parse_args()
    if not 0 <= args.overlay_alpha <= 255:
        parser.error("--overlay-alpha must be between 0 and 255.")
    with open(args.lyrics, encoding="utf-8") as file:
        lyrics_data = json.load(file)
    if not isinstance(lyrics_data, list) or not lyrics_data:
        parser.error("--lyrics must contain a non-empty JSON array.")
    create_pptx(lyrics_data, args.bg, args.output, parse_color(args.overlay_color), args.overlay_alpha, args.font_size, DELIVERY_QUALITY[args.delivery_mode], args.render_mode)


if __name__ == "__main__":
    main()
