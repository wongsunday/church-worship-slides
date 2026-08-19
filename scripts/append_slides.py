#!/usr/bin/env python3
"""Append slides while preserving a worship deck's image or editable render mode."""

import argparse
import json
import os
import sys

from pptx import Presentation

_here = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _here)
from render_slides import (  # noqa: E402
    DELIVERY_QUALITY,
    RENDER_MODES,
    add_editable_slide,
    auto_font_size,
    editable_auto_font_size,
    get_presentation_render_mode,
    parse_color,
    render_slide,
    set_presentation_render_mode,
    validate_background,
)


def main():
    parser = argparse.ArgumentParser(
        description="Append a new song to an image-rendered or editable worship PPTX"
    )
    parser.add_argument("--existing", required=True, help="Existing PPTX path")
    parser.add_argument("--lyrics", required=True, help="New song lyrics JSON path")
    parser.add_argument("--bg", required=True, help="New song background image path")
    parser.add_argument("--output", required=True, help="Output PPTX path; may equal --existing")
    parser.add_argument("--overlay-color", default="0,0,0")
    parser.add_argument("--overlay-alpha", type=int, default=185)
    parser.add_argument("--font-size", type=int, default=0)
    parser.add_argument("--delivery-mode", default="standard", choices=list(DELIVERY_QUALITY))
    parser.add_argument(
        "--render-mode",
        default=None,
        choices=RENDER_MODES,
        help="Required only for legacy v1.3.0-or-earlier decks without a saved mode marker",
    )
    args = parser.parse_args()
    if not 0 <= args.overlay_alpha <= 255:
        parser.error("--overlay-alpha must be between 0 and 255.")

    prs = Presentation(args.existing)
    recorded_mode = get_presentation_render_mode(prs)
    if recorded_mode and args.render_mode and args.render_mode != recorded_mode:
        parser.error(
            f"--render-mode {args.render_mode!r} conflicts with the existing deck's "
            f"{recorded_mode!r} mode. Regenerate the whole set to convert modes."
        )
    if recorded_mode is None and args.render_mode is None:
        parser.error(
            "This legacy deck has no render-mode marker. Pass --render-mode image "
            "to preserve its flattened format, or regenerate the whole set with "
            "--render-mode editable to create an editable deck."
        )
    render_mode = recorded_mode or args.render_mode
    existing_count = len(prs.slides)
    print(f"Existing PPTX: {existing_count} slide(s) | mode: {render_mode} | appending new song…")

    bg_img = validate_background(args.bg)
    with open(args.lyrics, encoding="utf-8") as file:
        lyrics_data = json.load(file)
    if not isinstance(lyrics_data, list) or not lyrics_data:
        parser.error("--lyrics must contain a non-empty JSON array.")

    all_lines = [line for slide in lyrics_data for line in slide["text"].split("\n")]
    font, used_size = auto_font_size(all_lines, args.font_size)
    if render_mode == "editable":
        used_size = editable_auto_font_size(all_lines, args.font_size)
    overlay_rgb = parse_color(args.overlay_color)
    preview_dir = os.path.join(os.path.dirname(os.path.abspath(args.output)), "_slide_imgs")
    os.makedirs(preview_dir, exist_ok=True)

    for slide_data in lyrics_data:
        file_index = existing_count + slide_data["index"]
        preview_path = os.path.join(preview_dir, f"slide_{file_index:02d}.jpg")
        render_slide(slide_data["text"], bg_img, font, overlay_rgb, args.overlay_alpha).save(
            preview_path, "JPEG", quality=DELIVERY_QUALITY[args.delivery_mode], optimize=True
        )
        if render_mode == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(preview_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            add_editable_slide(
                prs,
                slide_data["text"],
                args.bg,
                font,
                used_size,
                overlay_rgb,
                args.overlay_alpha,
            )
        print(f"  Appended slide {file_index:02d}: preview {os.path.getsize(preview_path) // 1024} KB")

    set_presentation_render_mode(prs, render_mode)
    prs.save(args.output)
    print(f"Saved: {args.output} ({len(prs.slides)} slides total, {os.path.getsize(args.output) // 1024} KB)")


if __name__ == "__main__":
    main()
